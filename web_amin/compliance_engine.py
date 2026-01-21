# compliance_engine.py - Version 3.1 CORRIGÃ‰E
import os
import json
import re
import httpx
import pandas as pd
from openai import OpenAI
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from docx import Document
from PyPDF2 import PdfReader
from dataclasses import dataclass
from typing import List, Dict, Tuple, Optional
import copy

from pptx_annotator import annotate_presentation

# ==========================================
# CONFIGURATION
# ==========================================
API_KEY = "sk-5312125f52c6491488218c98a934862a"
API_BASE = "https://tokenfactory.esprit.tn/api"
MODEL_NAME = "hosted_vllm/Llama-3.1-70B-Instruct"

http_client = httpx.Client(verify=False)
client = OpenAI(api_key=API_KEY, base_url=API_BASE, http_client=http_client)

# ==========================================
# KNOWLEDGE BASE
# ==========================================
class KnowledgeBase:
    def __init__(self, glossaire_path, registration_path, rules_json_path):
        self.glossaire_path = glossaire_path
        self.registration_path = registration_path
        self.rules_json_path = rules_json_path
        self.df_glossaire = None
        self.df_registration = None
        self.rulebook = {}
        self.load_data()

    def load_data(self):
        print("ðŸ“š Chargement de la Base de Connaissance...")
        if os.path.exists(self.rules_json_path):
            with open(self.rules_json_path, "r", encoding="utf-8") as f:
                raw_rules = json.load(f)
                # Convertir en format simple {id: text}
                self.rulebook = {
                    rid: rdata['text'] if isinstance(rdata, dict) else rdata 
                    for rid, rdata in raw_rules.items()
                }
            print(f"   âœ… Rulebook chargÃ© ({len(self.rulebook)} rÃ¨gles).")
        else:
            print(f"   âš ï¸ ALERTE : Fichier {self.rules_json_path} introuvable !")

        if os.path.exists(self.glossaire_path):
            try:
                self.df_glossaire = pd.read_excel(self.glossaire_path)
                print(f"   âœ… Glossaire chargÃ© ({len(self.df_glossaire)} lignes)")
            except Exception as e:
                print(f"   âš ï¸ Erreur lecture glossaire: {e}")

        if os.path.exists(self.registration_path):
            try:
                self.df_registration = pd.read_excel(self.registration_path)
                print(f"   âœ… Registration chargÃ© ({len(self.df_registration)} lignes)")
            except Exception as e:
                print(f"   âš ï¸ Erreur lecture registration: {e}")

    def get_authorized_countries(self, fund_name):
        if self.df_registration is None:
            return "DonnÃ©es non disponibles."
        try:
            mask = self.df_registration.iloc[:, 0].astype(str).str.contains(
                fund_name, case=False, na=False
            )
            fund_row = self.df_registration[mask]
            if fund_row.empty:
                return f"Fonds '{fund_name}' INTROUVABLE dans Registration."
            return fund_row.to_string(index=False)[:20000]
        except Exception as e:
            return f"Erreur lecture Registration: {str(e)}"

    def get_disclaimers(self, target_client):
        if self.df_glossaire is None:
            return "DonnÃ©es non disponibles."
        try:
            mask = self.df_glossaire.iloc[:, 0].astype(str).str.contains(
                str(target_client), case=False, na=False
            )
            disclaimer_string = self.df_glossaire[mask].to_string(index=False)
            return disclaimer_string[:20000] if disclaimer_string else self.df_glossaire.to_string(index=False)[:20000]
        except Exception as e:
            return f"Erreur lecture Glossaire: {str(e)}"



# ==========================================
# CLASSIFICATEUR HYBRIDE
# ==========================================
class HybridSlideClassifier:
    """Classification hybride combinant LLM et heuristiques pondÃ©rÃ©es"""
    
    CATEGORIES = [
        "COVER", "SLIDE_2", "FUND_PRESENTATION", "PERFORMANCE",
        "ESG", "TEAM", "RISKS", "VALUES", "MARKETING_GENERIC", "LEGAL_END"
    ]
    
    WEIGHTS = {"heuristic": 0.4, "llm": 0.6}
    
    def __init__(self, client: OpenAI):
        self.client = client
        
    def classify_slide(self, slide_text: str, slide_id: int, total_slides: int) -> Dict:
        """Classification avec fusion de scores"""
        heuristic_scores = self._heuristic_classification(slide_text, slide_id, total_slides)
        llm_category, llm_confidence = self._llm_classification(slide_text)
        
        final_category = self._merge_classifications(
            heuristic_scores, llm_category, llm_confidence
        )
        
        return {
            "type": final_category,
            "justification": f"Hybrid (H:{max(heuristic_scores, key=heuristic_scores.get)}, L:{llm_category})",
            "confidence": self._compute_confidence(heuristic_scores, llm_category, llm_confidence)
        }
    
    def _heuristic_classification(self, text: str, slide_id: int, total_slides: int) -> Dict[str, float]:
        """Calcul de scores heuristiques pour chaque catÃ©gorie"""
        scores = {cat: 0.0 for cat in self.CATEGORIES}
        lower = text.lower()
        word_count = len(text.split())
        
        # COVER (position critique)
        if slide_id == 1:
            scores["COVER"] += 0.5
        if word_count < 150:
            scores["COVER"] += 0.2
        if any(k in lower for k in ["oddo bhf", "asset management"]):
            scores["COVER"] += 0.3
        if re.search(r"\b(20\d{2}|janvier|fevrier|mars|april|mai|juin|juillet|aout|septembre|october|novembre|decembre)\b", lower):
            scores["COVER"] += 0.2
        
        # SLIDE_2
        if any(k in lower for k in ["profil de risque", "pays de commercialisation", "disclaimer standard"]):
            scores["SLIDE_2"] += 0.7
        if slide_id == 2:
            scores["SLIDE_2"] += 0.3
        
        # FUND_PRESENTATION
        strategy_keywords = ["strategie", "strategy", "allocation", "seuil", "ticket minimum", "investment policy"]
        scores["FUND_PRESENTATION"] += 0.15 * sum(1 for k in strategy_keywords if k in lower)
        
        # PERFORMANCE
        if re.search(r"\b\d{1,3}(?:[\.,]\d{1,2})?\s?%", text):
            scores["PERFORMANCE"] += 0.4
        if any(k in lower for k in ["ytd", "year to date", "annualise", "benchmark"]):
            scores["PERFORMANCE"] += 0.3
        
        # ESG
        if any(k in lower for k in ["esg", "sfdr", "sustainability", "article 8", "article 9"]):
            scores["ESG"] += 0.8
        
        # TEAM
        if any(k in lower for k in ["gerant", "manager", "team", "equipe de gestion"]):
            scores["TEAM"] += 0.6
        
        # RISKS
        risk_keywords = ["risque", "risk", "avertissement", "warning"]
        scores["RISKS"] += 0.2 * sum(1 for k in risk_keywords if k in lower)
        
        # LEGAL_END
        if slide_id >= total_slides - 2:
            scores["LEGAL_END"] += 0.4
        if any(k in lower for k in ["societe de gestion", "sgp", "r.c.s", "mention legale"]):
            scores["LEGAL_END"] += 0.4
        
        # VALUES
        if any(k in lower for k in ["valeur", "emetteur", "issuer", "stock"]):
            scores["VALUES"] += 0.5
        
        # Normalisation
        total = sum(scores.values())
        if total > 0:
            scores = {k: v/total for k, v in scores.items()}
        
        return scores
    
    def _llm_classification(self, text: str) -> Tuple[str, float]:
        """Classification via LLM avec extraction de confiance"""
        prompt = f"""Tu es un classificateur de slides.
CatÃ©gories: {', '.join(self.CATEGORIES)}

RÃ©ponds UNIQUEMENT en JSON:
{{
  "type": "<CATEGORIE>",
  "confidence": <0.0-1.0>,
  "justification": "<raison courte>"
}}

Texte de la slide:
{text[:1000]}
"""
        
        try:
            response = self.client.chat.completions.create(
                model=MODEL_NAME,
                messages=[
                    {"role": "system", "content": "Tu es un classificateur prÃ©cis et concis."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.0,
                max_tokens=150
            )
            
            raw = response.choices[0].message.content.strip()
            data = json.loads(re.search(r"\{.*\}", raw, re.DOTALL).group(0))
            
            category = data.get("type", "MARKETING_GENERIC")
            confidence = float(data.get("confidence", 0.5))
            
            if category in self.CATEGORIES:
                return category, confidence
            
        except Exception as e:
            print(f"   âš ï¸ Erreur LLM classification: {e}")
        
        return "MARKETING_GENERIC", 0.3
    
    def _merge_classifications(self, heuristic_scores: Dict, llm_category: str, llm_confidence: float) -> str:
        """Fusion intelligente des deux mÃ©thodes"""
        top_heuristic = max(heuristic_scores, key=heuristic_scores.get)
        heuristic_score = heuristic_scores[top_heuristic]
        
        weighted_heuristic = heuristic_score * self.WEIGHTS["heuristic"]
        weighted_llm = llm_confidence * self.WEIGHTS["llm"]
        
        if llm_confidence > 0.8:
            return llm_category
        
        if top_heuristic == llm_category:
            return top_heuristic
        
        return top_heuristic if weighted_heuristic > weighted_llm else llm_category
    
    def _compute_confidence(self, heuristic_scores: Dict, llm_category: str, llm_confidence: float) -> float:
        """Calcul de la confiance globale"""
        top_heuristic = max(heuristic_scores, key=heuristic_scores.get)
        
        if top_heuristic == llm_category:
            return (heuristic_scores[top_heuristic] + llm_confidence) / 2
        
        return max(heuristic_scores[top_heuristic], llm_confidence) * 0.7

# ==========================================
# RETRIEVER PROSPECTUS SÃ‰MANTIQUE
# ==========================================
@dataclass
class ProspectusChunk:
    """Fragment du prospectus avec mÃ©tadonnÃ©es"""
    text: str
    section: str
    page: int
    relevance_score: float = 0.0

class SemanticProspectusRetriever:
    """Extraction intelligente d'informations pertinentes du prospectus"""
    
    SECTION_KEYWORDS = {
        "STRATEGY": ["strategie", "objectif", "allocation", "investment policy"],
        "RISKS": ["risque", "profil de risque", "volatilite", "risk profile"],
        "PERFORMANCE": ["performance", "rendement", "historique", "benchmark"],
        "ESG": ["esg", "sfdr", "sustainability", "durabilite"],
        "LEGAL": ["prospectus", "reglement", "societe de gestion"],
        "FEES": ["frais", "commission", "costs", "fees"],
        "CHARACTERISTICS": ["caracteristiques", "isin", "part", "share class"]
    }
    
    def __init__(self, prospectus_text: str):
        self.full_text = prospectus_text
        self.chunks = self._chunk_prospectus(prospectus_text)
    
    def _chunk_prospectus(self, text: str, chunk_size: int = 800) -> List[ProspectusChunk]:
        """DÃ©coupe le prospectus en chunks avec dÃ©tection de sections"""
        chunks = []
        paragraphs = [p.strip() for p in text.split('\n') if p.strip()]
        
        current_section = "GENERAL"
        current_chunk = []
        current_length = 0
        page_counter = 1
        
        for para in paragraphs:
            detected_section = self._detect_section(para)
            if detected_section:
                current_section = detected_section
            
            current_chunk.append(para)
            current_length += len(para)
            
            if current_length >= chunk_size:
                chunks.append(ProspectusChunk(
                    text="\n".join(current_chunk),
                    section=current_section,
                    page=page_counter
                ))
                current_chunk = []
                current_length = 0
                page_counter += 1
        
        if current_chunk:
            chunks.append(ProspectusChunk(
                text="\n".join(current_chunk),
                section=current_section,
                page=page_counter
            ))
        
        return chunks
    
    def _detect_section(self, paragraph: str) -> Optional[str]:
        """DÃ©tecte le type de section Ã  partir d'un paragraphe"""
        lower = paragraph.lower()
        
        for section_name, keywords in self.SECTION_KEYWORDS.items():
            if any(keyword in lower for keyword in keywords):
                if len(paragraph.split()) < 15:
                    return section_name
        
        return None
    
    def retrieve_relevant_context(self, slide_text: str, slide_type: str, 
                                   top_k: int = 3, max_tokens: int = 2000) -> str:
        """RÃ©cupÃ¨re les chunks les plus pertinents pour une slide"""
        for chunk in self.chunks:
            chunk.relevance_score = self._compute_relevance(chunk, slide_text, slide_type)
        
        sorted_chunks = sorted(self.chunks, key=lambda c: c.relevance_score, reverse=True)
        selected_chunks = sorted_chunks[:top_k]
        
        context_parts = []
        current_length = 0
        
        for chunk in selected_chunks:
            chunk_text = f"[{chunk.section} - Page {chunk.page}]\n{chunk.text}"
            
            if current_length + len(chunk_text) > max_tokens:
                break
            
            context_parts.append(chunk_text)
            current_length += len(chunk_text)
        
        if not context_parts:
            return "CONTEXTE PROSPECTUS: Aucune information pertinente trouvÃ©e."
        
        return "\n\n---\n\n".join(context_parts)
    
    def _compute_relevance(self, chunk: ProspectusChunk, slide_text: str, slide_type: str) -> float:
        """Calcule un score de pertinence entre un chunk et une slide"""
        score = 0.0
        
        if slide_type in self.SECTION_KEYWORDS:
            section_keywords = self.SECTION_KEYWORDS[slide_type]
            for keyword in section_keywords:
                if keyword in chunk.section.lower():
                    score += 2.0
                    break
        
        slide_words = self._extract_keywords(slide_text)
        chunk_lower = chunk.text.lower()
        
        for word in slide_words:
            if word in chunk_lower:
                count = chunk_lower.count(word)
                score += count * 0.5
        
        rare_words = [w for w in slide_words if len(w) > 8]
        for word in rare_words:
            if word in chunk_lower:
                score += 1.5
        
        if len(chunk.text) < 200:
            score *= 0.7
        
        return score
    
    def _extract_keywords(self, text: str) -> List[str]:
        """Extraction de mots-clÃ©s significatifs"""
        text = re.sub(r'[^\w\s]', ' ', text.lower())
        words = text.split()
        
        stopwords = {'le', 'la', 'les', 'un', 'une', 'des', 'de', 'du', 'et', 'ou', 
                     'the', 'a', 'an', 'and', 'or', 'of', 'to', 'in', 'for'}
        
        keywords = [word for word in words if len(word) > 3 and word not in stopwords]
        
        return keywords

# ==========================================
# PROMPTS OPTIMISÃ‰S
# ==========================================
class OptimizedCompliancePrompts:
    """BibliothÃ¨que de prompts optimisÃ©s pour l'analyse de conformitÃ©"""
    
    @staticmethod
    def build_analysis_prompt(
        slide_id: int,
        total_slides: int,
        slide_type: str,
        slide_text: str,
        client_type: str,
        fund_name: str,
        rules_text: str,
        prospectus_context: str,
        max_rules: int = 10
    ) -> dict:
        """Construction d'un prompt structurÃ© avec exemples"""
        
        filtered_rules = OptimizedCompliancePrompts._filter_rules(rules_text, max_rules)
        
        # AJOUTER CETTE SECTION IMPORTANTE
        clarifications = """
âš ï¸ INTERPRÃ‰TATIONS IMPORTANTES DES RÃˆGLES :

1. RULE_018 (mois et annÃ©e) :
   - Format "September 2025" â†’ âœ… VALIDE (anglais)
   - Format "Septembre 2025" â†’ âœ… VALIDE (franÃ§ais)
   - Format "Sept 2025" â†’ âŒ NON VALIDE (abrÃ©viation)
   
2. RÃ¨gles de date : Accepter les formats mois-annÃ©e en franÃ§ais OU anglais
3. Performance : Ne signale pas si le disclaimer est prÃ©sent
"""
        
        system_prompt = f"""Tu es un expert en conformitÃ© rÃ©glementaire pour ODDO BHF Asset Management.

ðŸŽ¯ MISSION : Analyser la slide {slide_id}/{total_slides} et dÃ©tecter UNIQUEMENT les violations EXPLICITES.

{clarifications}

ðŸ“‹ CONTEXTE
â€¢ Type de slide: {slide_type}
â€¢ Client: {client_type}
â€¢ Fonds: {fund_name}

âš ï¸ RÃˆGLES STRICTES
1. Ne signale une violation QUE si tu trouves une PREUVE TEXTUELLE dans la slide
2. Cite EXACTEMENT le passage problÃ©matique (15 mots max)
3. Si la slide est conforme, rÃ©ponds {{"conforme": true, "violations": []}}
4. N'invente JAMAIS de violation
5. Sois PRAGMATIQUE : Ne signale pas des violations mineures ou subjectives

ðŸ“– RÃˆGLES APPLICABLES (SEULEMENT celles-ci)
{filtered_rules}

ðŸ“„ EXTRAIT PROSPECTUS PERTINENT
{prospectus_context[:1500]}

ðŸ” FORMAT DE RÃ‰PONSE (JSON strict)
{{
  "slide_id": {slide_id},
  "slide_type": "{slide_type}",
  "conforme": true|false,
  "violations": [
    {{
      "rule_id": "RULE_XXX",
      "rule_text": "Extrait de la rÃ¨gle",
      "issue": "Description prÃ©cise du problÃ¨me",
      "evidence": "Citation exacte de la slide (15 mots max)",
      "suggested_fix": "Action corrective CONCRÃˆTE"
    }}
  ],
  "notes": "Optionnel: explications supplÃ©mentaires"
}}
"""
        
        # ... reste inchangÃ©
        
        few_shot_examples = OptimizedCompliancePrompts._get_few_shot_examples(slide_type)
        
        user_message = f"""ðŸ“„ TEXTE DE LA SLIDE {slide_id}:
{slide_text[:1500]}

Analyse cette slide et rÃ©ponds en JSON. Si aucune violation TEXTUELLE, rÃ©ponds {{"conforme": true, "violations": []}}"""
        
        messages = [{"role": "system", "content": system_prompt}]
        
        for example in few_shot_examples:
            messages.append({"role": "user", "content": example["user"]})
            messages.append({"role": "assistant", "content": example["assistant"]})
        
        messages.append({"role": "user", "content": user_message})
        
        return {"messages": messages, "temperature": 0.05, "max_tokens": 800}
    
    @staticmethod
    def _filter_rules(rules_text: str, max_rules: int) -> str:
        """Filtre les rÃ¨gles pour ne garder que les plus pertinentes"""
        rules_lines = [line for line in rules_text.split('\n') if line.strip()]
        
        if len(rules_lines) > max_rules:
            rules_lines = rules_lines[:max_rules]
        
        return '\n'.join(rules_lines)
    
    @staticmethod
    def _get_few_shot_examples(slide_type: str) -> list:
        """Exemples few-shot adaptÃ©s au type de slide"""
        base_examples = [
            {
                "user": '''Slide 5: "Notre stratÃ©gie ESG permet d'investir dans des entreprises responsables. Performance YTD: +12.3%"
RÃ¨gles: RULE_090 - Afficher disclaimer performance''',
                "assistant": '''{
  "slide_id": 5,
  "conforme": false,
  "violations": [{
    "rule_id": "RULE_090",
    "rule_text": "Afficher disclaimer: 'les performances passÃ©es ne prÃ©sagent pas...'",
    "issue": "Performance affichÃ©e sans disclaimer obligatoire",
    "evidence": "Performance YTD: +12.3%",
    "suggested_fix": "Ajouter: 'Les performances passÃ©es ne prÃ©sagent pas des performances futures'"
  }]
}'''
            },
            {
                "user": '''Slide 3: "StratÃ©gie d'investissement: allocation flexible entre 60-80% actions europÃ©ennes. Ticket minimum: 1000â‚¬"
RÃ¨gles: RULE_008 - StratÃ©gie conforme au prospectus''',
                "assistant": '''{
  "slide_id": 3,
  "conforme": true,
  "violations": [],
  "notes": "StratÃ©gie correctement prÃ©sentÃ©e avec seuils d'investissement"
}'''
            }
        ]
        
        specific_examples = {
            "PERFORMANCE": [{
                "user": '''Slide 7: "Rendement annualisÃ© sur 5 ans: 8.2%"
RÃ¨gles: RULE_090 - Disclaimer performance''',
                "assistant": '''{
  "slide_id": 7,
  "conforme": false,
  "violations": [{
    "rule_id": "RULE_090",
    "issue": "Disclaimer de performance manquant",
    "evidence": "Rendement annualisÃ© sur 5 ans: 8.2%",
    "suggested_fix": "Ajouter disclaimer lisible: 'Les performances passÃ©es ne prÃ©sagent pas...'"
  }]
}'''
            }],
            "COVER": [{
                "user": '''Slide 1: "ODDO BHF Avenir Europe - Octobre 2024 - Document promotionnel - Client: Retail"
RÃ¨gles: RULE_018 - Page de garde doit indiquer nom, date, mention promotionnelle, cible''',
                "assistant": '''{
  "slide_id": 1,
  "conforme": true,
  "violations": [],
  "notes": "Page de garde complÃ¨te avec tous les Ã©lÃ©ments requis"
}'''
            }]
        }
        
        return base_examples + specific_examples.get(slide_type, [])

# ==========================================
# PARSERS (inchangÃ©s)
# ==========================================
def extract_text_from_pdf(pdf_path):
    if not os.path.exists(pdf_path):
        return ""
    try:
        reader = PdfReader(pdf_path)
        return "\n".join([page.extract_text() or "" for page in reader.pages])
    except Exception as e:
        print(f"âš ï¸ Erreur lecture PDF {pdf_path}: {e}")
        return ""

def extract_text_from_docx(docx_path):
    if not os.path.exists(docx_path):
        return ""
    try:
        doc = Document(docx_path)
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    except Exception as e:
        print(f"âš ï¸ Erreur lecture DOCX {docx_path}: {e}")
        return ""

def extract_prospectus_text(base_path):
    docx_file = os.path.join(base_path, "prospectus.docx")
    pdf_file = os.path.join(base_path, "prospectus.pdf")

    if os.path.exists(docx_file):
        print("ðŸ“„ Prospectus dÃ©tectÃ© : DOCX")
        return extract_text_from_docx(docx_file)

    if os.path.exists(pdf_file):
        print("ðŸ“„ Prospectus dÃ©tectÃ© : PDF")
        return extract_text_from_pdf(pdf_file)

    print("âŒ Aucun prospectus trouvÃ© (ni DOCX ni PDF)")
    return ""



# SUITE DE compliance_engine.py

# Parsers PPTX
def _get_text_from_text_frame(text_frame):
    parts = []
    for para in text_frame.paragraphs:
        run_texts = [run.text for run in para.runs if run.text]
        if run_texts:
            parts.append("".join(run_texts))
        else:
            if para.text:
                parts.append(para.text)
    return "\n".join(parts).strip()

def _collect_shape_text(shape):
    try:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            texts = []
            for sh in shape.shapes:
                t = _collect_shape_text(sh)
                if t:
                    texts.append(t)
            return "\n".join(t for t in texts if t).strip()

        if hasattr(shape, "text") and isinstance(shape.text, str) and shape.text.strip():
            return shape.text.replace('\x0b', ' ').strip()

        if hasattr(shape, "text_frame") and shape.text_frame is not None:
            tf_text = _get_text_from_text_frame(shape.text_frame)
            if tf_text:
                return tf_text

        if getattr(shape, "is_placeholder", False):
            try:
                if hasattr(shape, "text_frame") and shape.text_frame is not None:
                    tf_text = _get_text_from_text_frame(shape.text_frame)
                    if tf_text:
                        return tf_text
            except Exception:
                pass

        if hasattr(shape, "alternative_text") and shape.alternative_text:
            return shape.alternative_text.strip()

    except Exception:
        return ""

    return ""

def extract_slides_data(pptx_path):
    if not os.path.exists(pptx_path):
        return []

    prs = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            t = _collect_shape_text(shape)
            if t:
                texts.append(t)

        text = "\n".join([t for t in texts if t]).strip()
        slides.append({"id": i + 1, "text": text})
    return slides

def load_metadata(folder_path):
    with open(os.path.join(folder_path, 'metadata.json'), 'r', encoding='utf-8') as f:
        raw = json.load(f)
    
    metadata = {}
    metadata["management_company"] = raw.get("SociÃ©tÃ© de Gestion", "")
    metadata["is_sicav"] = raw.get("Est ce que le produit fait partie de la Sicav d'Oddo", False)
    
    prof = raw.get("Le client est-il un professionnel", False)
    if isinstance(prof, str):
        prof = prof.strip().lower() in ["true", "yes", "1", "oui"]
    metadata["client_is_professional"] = bool(prof)
    
    metadata["new_strategy"] = raw.get("Le document fait-il rÃ©fÃ©rence Ã  une nouvelle StratÃ©gie", False)
    metadata["new_product"] = raw.get("Le document fait-il rÃ©fÃ©rence Ã  un nouveau Produit", False)
    metadata["fund_name"] = raw.get("fund_name", "")
    metadata["client_type"] = "Professional" if metadata["client_is_professional"] else "Retail"
    
    return metadata

# ==========================================
# RULES MANAGEMENT
# ==========================================
def build_auto_rules_mapping(rulebook):
    """Mapping amÃ©liorÃ© avec gestion des rÃ¨gles enrichies"""
    mapping = {cat: [] for cat in HybridSlideClassifier.CATEGORIES}
    
    cat_keywords = {
        "COVER": ["page de garde", "document promotionnel", "mois"],
        "SLIDE_2": ["profil de risque", "pays de commercialisation", "slide 2"],
        "FUND_PRESENTATION": ["strategie du fonds", "allocation geographique"],
        "PERFORMANCE": ["performance", "ytd", "benchmark", "rendement"],
        "ESG": ["esg", "sfdr"],
        "TEAM": ["equipe", "gerants"],
        "RISKS": ["risque", "risk", "avertissement"],
        "VALUES": ["valeur", "emetteur"],
        "MARKETING_GENERIC": ["avantages", "why"],
        "LEGAL_END": ["mention legale", "sgp"]
    }

    for rule_id, rule_text in rulebook.items():
        # Si c'est un dict (rÃ¨gle enrichie), extraire le texte
        if isinstance(rule_text, dict):
            rule_text = rule_text.get('text', '')
        
        rtxt = rule_text.lower()
        assigned = False
        
        for cat, kws in cat_keywords.items():
            for kw in kws:
                if kw in rtxt:
                    mapping[cat].append(rule_id)
                    assigned = True
                    break
            if assigned:
                break
        
        if not assigned:
            mapping["MARKETING_GENERIC"].append(rule_id)

    for k in mapping:
        mapping[k] = sorted(list(dict.fromkeys(mapping[k])))
    
    return mapping

def build_rule_conditions(rulebook):
    """Conditions intelligentes et contextuelles"""
    conditions = {}
    
    for rid, rule_data in rulebook.items():
        # Extraire le texte (compatibilitÃ© rÃ¨gles enrichies)
        txt = rule_data if isinstance(rule_data, str) else rule_data.get('text', '')
        ltxt = txt.lower()
        
        def _always(metadata, slide_type, slide_id=None, total_slides=None):
            return True
        cond = _always

        # Professionnels uniquement
        if any(k in ltxt for k in ["professionnel", "professional", "clients professionnels"]):
            cond = (lambda _prev: (lambda metadata, slide_type, slide_id=None, total_slides=None: 
                metadata.get("client_is_professional", False)))(_always)

        # Retail uniquement
        elif any(k in ltxt for k in ["retail", "client non professionnel", "presentations retail"]):
            cond = (lambda _prev: (lambda metadata, slide_type, slide_id=None, total_slides=None: 
                not metadata.get("client_is_professional", False)))(_always)

        # PrÃ©-commercialisation
        if "pre-commercialisation" in ltxt or "pre commercialisation" in ltxt:
            cond = (lambda prev: (lambda metadata, slide_type, slide_id=None, total_slides=None: 
                metadata.get("new_product", False) and prev(metadata, slide_type, slide_id, total_slides)))(cond)

        # Registration (ne s'applique PAS aux nouveaux produits)
        if "registration abroad" in ltxt or "pays de commercialisation" in ltxt:
            cond = (lambda prev: (lambda metadata, slide_type, slide_id=None, total_slides=None: 
                (not metadata.get("new_product", False)) and prev(metadata, slide_type, slide_id, total_slides)))(cond)

        # COVER
        if any(k in ltxt for k in ["page de garde", "doit indiquer"]) and "commencer" not in ltxt:
            cond = (lambda prev: (lambda metadata, slide_type, slide_id=None, total_slides=None: 
                slide_type == "COVER" and prev(metadata, slide_type, slide_id, total_slides)))(cond)

        # LEGAL_END
        if any(k in ltxt for k in ["mention legale de la sgp", "mention legale"]) and "page de garde" not in ltxt:
            cond = (lambda prev: (lambda metadata, slide_type, slide_id=None, total_slides=None: 
                slide_type == "LEGAL_END" and prev(metadata, slide_type, slide_id, total_slides)))(cond)

        # StratÃ©gie
        if "strategie du fonds" in ltxt:
            cond = (lambda prev: (lambda metadata, slide_type, slide_id=None, total_slides=None: 
                slide_type == "FUND_PRESENTATION" and prev(metadata, slide_type, slide_id, total_slides)))(cond)

        # Performance
        if any(k in ltxt for k in ["performances passees", "past performance", "simulation"]):
            cond = (lambda prev: (lambda metadata, slide_type, slide_id=None, total_slides=None: 
                slide_type in ["PERFORMANCE", "FUND_PRESENTATION"] and prev(metadata, slide_type, slide_id, total_slides)))(cond)

        conditions[rid] = cond
    
    return conditions

def extract_rules_subset_text(rulebook, active_rules_ids, rule_conditions, metadata, slide_type, slide_id=None, total_slides=None):
    """Extraction avec gestion des rÃ¨gles enrichies"""
    lines = []
    
    for rid in active_rules_ids:
        cond = rule_conditions.get(rid, lambda m, s, sid=None, ts=None: True)
        
        try:
            if cond(metadata, slide_type, slide_id, total_slides):
                # Extraire le texte (compatibilitÃ©)
                rule_data = rulebook.get(rid, "")
                txt = rule_data if isinstance(rule_data, str) else rule_data.get('text', '')
                
                if txt:
                    lines.append(f"{rid}: {txt.strip()}")
        except Exception as e:
            print(f"   âš ï¸ Erreur condition {rid}: {e}")
            continue
    
    return "\n".join(lines) if lines else "Aucune rÃ¨gle applicable extraite."

# ==========================================
# DÃ‰TECTIONS DÃ‰TERMINISTES CORRIGÃ‰ES
# ==========================================
RE_RECOMMEND = re.compile(r"\b(acheter|vendre|renforcer|buy|sell)\b", re.I)
RE_PERF_TOKENS = re.compile(r"\b(ytd|annualis|performance|rendement|benchmark)\b", re.I)
RE_PERCENT = re.compile(r"\b\d{1,3}(?:[\.,]\d{1,2})?\s?%")
RE_DISCLAIMER_PERF = re.compile(r"(les performances passees ne|past performance|not a reliable indicator|ne presagent pas)", re.I)
RE_HISTORICAL_DATA = re.compile(r"'[89]\d|19[89]\d|20[012]\d|historique.*\d+\s*ans", re.I)
RE_SHARE_CLASS = re.compile(r"\b(CR|DR|CN|DN|GC|part\s+[A-Z]|classe\s+[A-Z]|share\s+class)\b", re.I)
RE_AFFIRMATIVE_MARKETING = re.compile(r"\b(why (invest|choose)|key component|must invest|should invest|you (should|must|need to)|il faut|vous devez)\b", re.I)
RE_ATTENUATION = re.compile(r"\b(selon (nous|notre|nos)|we (believe|think)|peut etre|pourrait|in our (view|opinion)|may|might|could)\b", re.I)

def deterministic_checks(slide_text, metadata, kb, slide_type, active_rules_ids, rule_conditions, slide_id, total_slides):
    """DÃ©tections CORRIGÃ‰ES plus intelligentes avec contexte"""
    v = []
    txt = slide_text
    lower = txt.lower()

    # 1) Recommendation language (RULE_049 dans les nouvelles rÃ¨gles)
    recommend_rule_id = "RULE_049" if "RULE_049" in active_rules_ids else None
    
    if recommend_rule_id and rule_conditions.get(recommend_rule_id, lambda m, s, sid, ts: True)(metadata, slide_type, slide_id, total_slides):
        # Chercher recommandations directes CLAIRES
        m = RE_RECOMMEND.search(txt)
        if m:
            v.append({
                "rule_id": recommend_rule_id,
                "rule_text": kb.rulebook.get(recommend_rule_id, "L'usage de formules ..."),
                "issue": "Formulation pouvant constituer une recommandation d'investissement.",
                "evidence": m.group(0),
                "suggested_fix": "Retirer ou reformuler avec attÃ©nuation ('selon nous', 'peut Ãªtre considÃ©rÃ©')."
            })
            return v, True
        
        # Formulations marketing trop affirmatives
        affirmative_match = RE_AFFIRMATIVE_MARKETING.search(txt)
        has_attenuation = RE_ATTENUATION.search(txt)
        has_disclaimer = re.search(r"(no investment recommendation|ne constitue pas une recommandation)", txt, re.I)
        
        if affirmative_match and not has_attenuation and not has_disclaimer:
            v.append({
                "rule_id": recommend_rule_id,
                "rule_text": kb.rulebook.get(recommend_rule_id, "L'usage de formules ..."),
                "issue": "Formulation marketing trop affirmative sans attÃ©nuation appropriÃ©e.",
                "evidence": f"'{affirmative_match.group(0)}' sans formule d'attÃ©nuation",
                "suggested_fix": "Ajouter des formules d'attÃ©nuation : 'selon notre analyse', 'peut Ãªtre considÃ©rÃ© comme', 'in our view'."
            })

    # 2) Performance detection
    perf_found = bool(RE_PERCENT.search(txt) or RE_PERF_TOKENS.search(txt))
    disclaimer_found = bool(RE_DISCLAIMER_PERF.search(txt))
    
    if perf_found:
        # VÃ©rifier historique
        has_long_history = bool(RE_HISTORICAL_DATA.search(txt))
        
        # RULE_064 : Historique < 1 an
        if "RULE_064" in active_rules_ids and not has_long_history:
            if rule_conditions.get("RULE_064", lambda m, s, sid, ts: True)(metadata, slide_type, slide_id, total_slides):
                year_matches = re.findall(r"'?\d{2}\s|19\d{2}|20[012]\d", txt)
                if len(year_matches) < 2:
                    v.append({
                        "rule_id": "RULE_064",
                        "rule_text": kb.rulebook.get("RULE_064", "Historique < 1 an..."),
                        "issue": "Impossible de vÃ©rifier si le fonds a plus d'un an d'historique.",
                        "evidence": "Aucune donnÃ©e temporelle claire dans la slide",
                        "suggested_fix": "PrÃ©ciser explicitement la date de crÃ©ation du fonds ou retirer les performances si < 1 an."
                    })
        
        # RULE_090 : Disclaimer performance manquant (CORRIGÃ‰)
        if not disclaimer_found:
            perf_rule = "RULE_090" if "RULE_090" in active_rules_ids else None
            
            if perf_rule and rule_conditions.get(perf_rule, lambda m, s, sid, ts: True)(metadata, slide_type, slide_id, total_slides):
                v.append({
                    "rule_id": perf_rule,
                    "rule_text": kb.rulebook.get(perf_rule, "Afficher le disclaimer de performance requis."),
                    "issue": "Performance affichÃ©e sans disclaimer obligatoire.",
                    "evidence": "PrÃ©sence de chiffres/termes de performance sans disclaimer",
                    "suggested_fix": "Afficher le disclaimer obligatoire 'les performances passÃ©es ne prÃ©sagent pas...' de maniÃ¨re lisible."
                })
                return v, True
        
        # RULE_057 : Parts retail
        if not metadata.get("client_is_professional", False):
            if "RULE_057" in active_rules_ids:
                if rule_conditions.get("RULE_057", lambda m, s, sid, ts: True)(metadata, slide_type, slide_id, total_slides):
                    if not RE_SHARE_CLASS.search(txt):
                        v.append({
                            "rule_id": "RULE_057",
                            "rule_text": kb.rulebook.get("RULE_057", "Parts retail uniquement..."),
                            "issue": "Pour un client retail, la classe de parts doit Ãªtre prÃ©cisÃ©e (CR, DR, CN, DN, GC).",
                            "evidence": "Aucune mention de classe de parts dÃ©tectÃ©e",
                            "suggested_fix": "PrÃ©ciser explicitement la classe de parts retail affichÃ©e."
                        })

    return v, False

# ==========================================
# ANALYSE AVEC CONTEXTE (VERSION FINALE)
# ==========================================
def analyze_compliance(slide, metadata, prospectus_retriever, kb, rules_mapping, rule_conditions, total_slides):
    """Analyse CORRIGÃ‰E avec contexte complet"""
    fund_name = metadata.get('fund_name', '')
    client_type = metadata.get('client_type', 'Retail')
    slide_id = slide['id']

    # Classification
    classifier = HybridSlideClassifier(client)
    slide_type_info = classifier.classify_slide(slide['text'], slide_id, total_slides)
    slide_type = slide_type_info.get("type", "MARKETING_GENERIC")

    # Contexte prospectus
    relevant_context = prospectus_retriever.retrieve_relevant_context(
        slide['text'], slide_type, top_k=3, max_tokens=2000
    )

    active_rules_ids = rules_mapping.get(slide_type, [])

    # CORRECTION : Filtrer RULE_025 si ce n'est pas la slide 1
    if "RULE_025" in active_rules_ids and slide_id != 1:
        active_rules_ids = [r for r in active_rules_ids if r != "RULE_025"]

    # DÃ©tections dÃ©terministes
    deterministic_violations, deterministic_only = deterministic_checks(
        slide['text'], metadata, kb, slide_type, active_rules_ids, rule_conditions, slide_id, total_slides
    )
    
    if deterministic_only and deterministic_violations:
        return {
            "slide_id": slide_id,
            "slide_type": slide_type,
            "conforme": False,
            "violations": deterministic_violations,
            "notes": "Violations dÃ©tectÃ©es par rÃ¨gles dÃ©terministes."
        }

    # Extraction rÃ¨gles applicables
    rules_corpus_text = extract_rules_subset_text(kb.rulebook, active_rules_ids, rule_conditions, metadata, slide_type, slide_id, total_slides)

    # Prompt optimisÃ©
    prompt_config = OptimizedCompliancePrompts.build_analysis_prompt(
        slide_id=slide_id,
        total_slides=total_slides,
        slide_type=slide_type,
        slide_text=slide['text'],
        client_type=client_type,
        fund_name=fund_name,
        rules_text=rules_corpus_text,
        prospectus_context=relevant_context,
        max_rules=10
    )

    try:
        resp = client.chat.completions.create(model=MODEL_NAME, **prompt_config)
        raw = resp.choices[0].message.content.strip()
        
        try:
            parsed = json.loads(raw)
            return merge_with_deterministic(parsed, deterministic_violations)
        except Exception:
            m = re.search(r"(\{(?:.|\n)*\})", raw)
            if m:
                try:
                    parsed = json.loads(m.group(1))
                    return merge_with_deterministic(parsed, deterministic_violations)
                except Exception:
                    pass
            
            return {
                "slide_id": slide_id,
                "slide_type": slide_type,
                "conforme": False,
                "violations": deterministic_violations,
                "notes": "Impossible de parser la rÃ©ponse LLM."
            }
    except Exception as e:
        return {
            "slide_id": slide_id,
            "slide_type": slide_type,
            "conforme": False,
            "violations": deterministic_violations,
            "notes": f"Erreur appel LLM: {e}"
        }

def merge_with_deterministic(parsed, deterministic_violations):
    """Fusion des violations LLM et dÃ©terministes"""
    if not deterministic_violations:
        return parsed
    
    parsed.setdefault("violations", [])
    existing_ids = {v.get("rule_id") for v in parsed.get("violations", [])}
    
    for dv in deterministic_violations:
        if dv.get("rule_id") not in existing_ids:
            parsed["violations"].append(dv)
            parsed["conforme"] = False
    
    return parsed

# ==========================================
# FONCTION PRINCIPALE MISE À JOUR
# ==========================================
def run_compliance_analysis(folder_path):
    """Fonction principale appelée par Flask"""
    print(f"🚀 Analyse de conformité démarrée")
    
    # Initialisation de la base de connaissance
    # Assurez-vous que les fichiers Excel/JSON sont bien dans le dossier racine du projet
    kb = KnowledgeBase(
        glossaire_path="GLOSSAIRE DISCLAIMERS 20231122 .xlsx",
        registration_path="Registration abroad of Funds_20251008.xlsx",
        rules_json_path="rules_enhanced.json" 
    )

    rules_mapping = build_auto_rules_mapping(kb.rulebook)
    rule_conditions = build_rule_conditions(kb.rulebook)

    try:
        metadata = load_metadata(folder_path)
        pptx_file = os.path.join(folder_path, 'presentation.pptx')
        
        # 1. Extraction et Analyse
        slides = extract_slides_data(pptx_file)
        prospectus_full_text = extract_prospectus_text(folder_path)
        total_slides = len(slides)

        prospectus_retriever = SemanticProspectusRetriever(prospectus_full_text)

        results = []
        for slide in slides:
            print(f" > Analyse Slide {slide['id']}...")
            res = analyze_compliance(
                slide, metadata, prospectus_retriever, 
                kb, rules_mapping, rule_conditions, total_slides
            )
            results.append(res)
        
        # 2. Génération PPTX annoté avec le NOUVEAU module
        # On définit le nom exact attendu par app.py
        annotated_path = os.path.join(folder_path, 'presentation_annotated.pptx')
        
        print(f"🎨 Lancement de l'annotateur avancé (COM)...")
        try:
            # Appel à votre nouveau script pptx_annotator
            annotate_presentation(pptx_file, results, output_path=annotated_path)
            
            if os.path.exists(annotated_path):
                print("✅ Annotation COM terminée avec succès.")
            else:
                print("⚠️ L'annotateur n'a pas généré le fichier.")
                
        except Exception as e_annot:
            print(f"❌ Erreur lors de l'annotation PowerPoint : {e_annot}")
            print("💡 Conseil : Vérifiez que PowerPoint est installé et qu'aucun processus n'est bloqué.")

        return results

    except Exception as e:
        print(f"❌ Erreur critique dans l'analyse: {e}")
        import traceback
        traceback.print_exc()
        raise