# ceci le le metadata enriched celle de pptx


# ═══════════════════════════════════════════════════════════════════
# CODE FINAL PPTX LAYER 1.5 BIS AMÉLIORÉ : EXTRACTION ROBUSTE + LLM
# + TABLE "GENERAL CHARACTERISTICS" (INCEPTION DATE + COUNTRIES, ETC.)
# ═══════════════════════════════════════════════════════════════════



import re
import json
from pathlib import Path
from typing import Dict, List, Any, Optional
from datetime import datetime

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import logging

# LLM
from openai import OpenAI
import httpx

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class ImprovedPPTXExtractor:
    """
    Extracteur PPTX amélioré avec :
      - Patterns robustes (regex)
      - Parsing structuré de la table "General characteristics"
      - LLM (Llama 3.1 70B) pour enrichissement facultatif

    Objectif : obtenir un fund_info_complete riche et STABLE, même si le LLM timeout.
    """

    def __init__(self, pptx_path: str, verbose: bool = True):
        self.pptx_path = pptx_path
        self.verbose = verbose
        self.prs = Presentation(pptx_path)
        self.filename = Path(pptx_path).name

        # Cache du texte complet
        self.full_text = self._extract_all_text()
        self.slides_data = self._extract_slides_with_metadata()

        # Cache pour la table "General characteristics"
        self._general_char_table: Optional[Dict[str, str]] = None

        self._log(f"📂 Chargement PPTX: {self.filename}")
        self._log(f"📊 Nombre de slides: {len(self.prs.slides)}")
        self._log(f"📝 Texte total: {len(self.full_text)} caractères\n")

    # ─────────────────────────────────────────────────────────────
    # NOUVELLE MÉTHODE : EXTRACTION COMPLÈTE AVEC LLM
    # ─────────────────────────────────────────────────────────────
    def extract_complete_metadata_with_llm(
        self,
        api_key: str,
        base_url: str = "https://tokenfactory.esprit.tn/api",
        model: str = "hosted_vllm/Llama-3.1-70B-Instruct"
    ) -> Dict[str, Any]:
        """
        Extraction complète avec LLM pour obtenir TOUTES les métadonnées importantes.

        Pipeline :
          1) Extraction déterministe (regex + table General characteristics)
          2) Localisation du slide "Fund / General characteristics"
          3) Appel LLM pour compléter / affiner certains champs
          4) Fusion robuste (patterns = socle, LLM = enrichissement)
        """
        self._log("\n🧠 EXTRACTION COMPLÈTE AVEC LLM")
        self._log("─" * 70 + "\n")

        # 1. Extraction de base (patterns + table)
        basic_extraction = self.extract_all()

        # 2. Trouver le slide "Fund / General characteristics"
        fund_char_slide = self._find_fund_characteristics_slide()

        if fund_char_slide:
            fund_char_text = self._get_slide_text(fund_char_slide)
            self._log(f"✅ Slide 'Fund/General characteristics' trouvé\n")
        else:
            fund_char_text = ""
            self._log(
                "⚠️ Slide 'Fund/General characteristics' non trouvé, "
                "utilisation du texte complet\n"
            )

        # 3. Appel LLM pour extraction structurée
        llm_metadata = self._call_llm_for_metadata_extraction(
            fund_char_text if fund_char_text else self.full_text[:15000],
            basic_extraction,
            api_key,
            base_url,
            model
        )

        # 4. Fusion des résultats
        complete_metadata = self._merge_extractions(basic_extraction, llm_metadata)

        return complete_metadata

    # ─────────────────────────────────────────────────────────────
    # LOCALISATION DU SLIDE "GENERAL / FUND CHARACTERISTICS"
    # ─────────────────────────────────────────────────────────────
    def _find_fund_characteristics_slide(self) -> Optional[Any]:
        """Trouve le slide contenant les caractéristiques du fonds."""
        keywords = [
            "fund characteristic",
            "general characteristic",
            "general characteristics",
            "structure and technical information",
            "fund details",
            "détails du fonds",
            "caractéristiques du fonds",
        ]

        for slide in self.prs.slides:
            text = self._get_slide_text(slide).lower()
            if any(kw in text for kw in keywords):
                return slide

        return None

    # ─────────────────────────────────────────────────────────────
    # PARSING STRUCTURÉ DE LA TABLE "GENERAL CHARACTERISTICS"
    # ─────────────────────────────────────────────────────────────
    def _normalize_label(self, label: str) -> str:
        """Normalise un libellé de table (minuscules, espaces homogènes)."""
        label = re.sub(r"[:\u00a0]+", " ", label)  # supprimer ':' et nbsp
        label = re.sub(r"\s+", " ", label).strip().lower()
        return label

    def _extract_general_characteristics_table(self) -> Dict[str, str]:
        """
        Parse la table 'General characteristics' si elle existe.

        Retourne un dict clé métier → valeur brute (string), par ex :
          {
            'countries_available_for_sales': 'Ireland, France, Germany, Switzerland',
            'inception_date': '2025',
            'domicile': 'Ireland',
            ...
          }
        """
        if self._general_char_table is not None:
            return self._general_char_table

        table_info: Dict[str, str] = {}

        # Localiser d'abord le slide pertinent
        slide = self._find_fund_characteristics_slide()
        if not slide:
            self._general_char_table = {}
            return self._general_char_table

        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
                continue

            table = shape.table
            for row in table.rows:
                cells = row.cells
                if len(cells) < 2:
                    continue

                raw_label = cells[0].text.strip()
                raw_value = cells[1].text.strip()

                if not raw_label or not raw_value:
                    continue

                label = self._normalize_label(raw_label)

                # Mapping label → clé canonique
                key = None
                if "countries available for sales" in label:
                    key = "countries_available_for_sales"
                elif "inception date" in label:
                    key = "inception_date"
                elif label.startswith("isin"):
                    key = "isin_code"
                elif label.startswith("domicile"):
                    key = "domicile"
                elif "management company" in label:
                    key = "management_company"
                elif "legal structure" in label:
                    key = "legal_structure"
                elif "stock exchange" in label:
                    key = "stock_exchanges"
                elif "tickers" in label or "ticker" in label:
                    key = "tickers"
                elif "dividend policy" in label:
                    key = "dividend_policy"
                elif "minimum initial investment amount" in label:
                    key = "minimum_initial_investment_amount"
                elif "creation unit" in label:
                    key = "creation_unit"
                elif "initial offering price" in label:
                    key = "initial_offering_price"
                elif label.startswith("dealing day"):
                    key = "dealing_day"
                elif label.startswith("dealing deadline"):
                    key = "dealing_deadline"
                elif "settlement date" in label:
                    key = "settlement_date_subscriptions"

                if key:
                    table_info[key] = raw_value.strip()

        self._general_char_table = table_info
        return self._general_char_table

    def _parse_countries_list(self, raw: str) -> List[str]:
        """Transforme 'Ireland, France, Germany, Switzerland' → ['Ireland', ...]."""
        if not raw:
            return []
        # Séparer par virgule, point-virgule ou slash
        parts = re.split(r"[;,/]", raw)
        countries = [p.strip() for p in parts if p.strip()]
        return countries

    # ─────────────────────────────────────────────────────────────
    # APPEL LLM POUR MÉTADONNÉES COMPLÉMENTAIRES
    # ─────────────────────────────────────────────────────────────
    def _call_llm_for_metadata_extraction(
        self,
        fund_text: str,
        basic_extraction: Dict[str, Any],
        api_key: str,
        base_url: str,
        model: str,
    ) -> Dict[str, Any]:
        """
        Appelle le LLM pour extraire les métadonnées de manière structurée.
        Utilisé en complément des extractions déterministes.
        """
        self._log("🔄 Appel du LLM pour extraction structurée...")

        http_client = httpx.Client(verify=False)
        client = OpenAI(api_key=api_key, base_url=base_url, http_client=http_client)

        prompt = f"""You are a financial document analyzer specialized in fund documentation.

Extract ALL important metadata from this fund presentation text.

TEXT TO ANALYZE:
<<<
{fund_text}
>>>

ALREADY EXTRACTED (basic patterns + table parsing):
{json.dumps(basic_extraction['fund_info_complete'], indent=2, ensure_ascii=False)}

Return a JSON object with these fields (extract ONLY if clearly present, otherwise null):

{{
  "countries_available_for_sales": ["country1", "country2", ...],
  "isin_code": "string or null",
  "domicile": "string or null",
  "management_company": "string or null",
  "legal_structure": "string or null",
  "custodian": "string or null",
  "stock_exchanges": ["exchange1", "exchange2", ...],
  "tickers": ["ticker1", "ticker2", ...],
  "creation_unit": "string or null",
  "investment_manager": "string or null",
  "recommended_holding_period": "string or null",
  "dealing_deadline": "string or null",
  "settlement_date": "string or null",
  "fund_size": "string or null",
  "aum": "string or null"
}}

IMPORTANT for countries_available_for_sales:
- Look for phrases like: "Countries in which the fund is authorized", "available for distribution", "licensed for sale".
- Typical countries: Ireland, Germany, France, Switzerland, Belgium, Luxembourg, etc.

Return ONLY valid JSON, no markdown, no comments."""

        try:
            response = client.chat.completions.create(
                model=model,
                messages=[
                    {
                        "role": "system",
                        "content": "You are a precise JSON extractor. Return only valid JSON.",
                    },
                    {"role": "user", "content": prompt},
                ],
                temperature=0.0,
                max_tokens=900,
            )

            content = response.choices[0].message.content.strip()

            # Nettoyage éventuel de ```json
            if "```json" in content:
                content = content.split("```json")[1].split("```")[0]
            elif "```" in content:
                content = content.split("```")[1].split("```")[0]

            data = json.loads(content.strip())

            self._log("   ✅ LLM : extraction réussie\n")
            self._log("   📋 Métadonnées extraites par LLM:")
            for key, value in data.items():
                if value and value != "null":
                    self._log(f"      • {key}: {value}")

            return data

        except Exception as e:
            self._log(f"   ❌ Erreur LLM (timeout ou autre): {e}")
            # En cas d'échec LLM, on ne casse pas la pipeline
            return {}

    # ─────────────────────────────────────────────────────────────
    # FUSION DES EXTRACTIONS
    # ─────────────────────────────────────────────────────────────
    def _merge_extractions(
        self, basic: Dict[str, Any], llm_data: Dict[str, Any]
    ) -> Dict[str, Any]:
        """
        Fusionne les extractions de base et LLM.

        Règle métier :
          - Si basic a une valeur non nulle → on la garde (source déterministe)
          - Sinon → on prend la valeur LLM si disponible
        """
        self._log("\n🔄 FUSION DES EXTRACTIONS")
        self._log("─" * 70 + "\n")

        merged = basic["fund_info_complete"].copy()

        for key, value in llm_data.items():
            if key not in merged or merged.get(key) in [None, "", "None", [], "null"]:
                merged[key] = value if value not in ["null", ""] else None
                symbol = "✅"
            else:
                symbol = "🔒"  # valeur de base conservée

            value_str = str(value)[:60] if value else "Non trouvé"
            self._log(f"{symbol} {key:<35} {value_str}")

        result = {
            "fund_info_complete": merged,
            "fund_info_detailed": basic["fund_info_detailed"],
            "extraction_sources": basic["extraction_sources"],
            "performance_flags": basic["performance_flags"],
            "llm_metadata": llm_data,
            "extraction_metadata": {
                "extraction_date": datetime.now().isoformat(),
                "source_file": self.filename,
                "total_slides": len(self.prs.slides),
                "extraction_method": "improved_with_llm",
            },
        }

        return result

    # ─────────────────────────────────────────────────────────────
    # MÉTHODE PRINCIPALE : EXTRACTION DE BASE (ROBUSTE)
    # ─────────────────────────────────────────────────────────────
    def extract_all(self) -> Dict[str, Any]:
        """Extrait toutes les données avec méthodes améliorées (patterns + table)."""

        self._log("🔍 EXTRACTION DONNÉES PPTX (PATTERNS + TABLE)")
        self._log("─" * 70 + "\n")

        # 1. Parser la table "General characteristics" une seule fois
        table_info = self._extract_general_characteristics_table()

        # 2. Extraction principale
        inception_from_table = table_info.get("inception_date")
        inception_date = self._extract_inception_date_robust(
            inception_from_table=inception_from_table
        )

        extracted: Dict[str, Any] = {
            "inception_date": inception_date,
            "sri_risk_level": self._extract_sri_risk_level_robust(),
            "esg_sfdr_article": self._extract_esg_sfdr_article_robust(),
            "benchmark": self._extract_benchmark_robust(),
            "fund_type": self._extract_fund_type_robust(),
            "minimum_investment": self._extract_minimum_investment_robust(),
            "ter": self._extract_ter_robust(),
            "distribution_policy": self._extract_distribution_policy_robust(),
            "currency": self._extract_currency_robust(),
            "fund_name": self._extract_fund_name(),
            # Champs issus de la table "General characteristics"
            "countries_available_for_sales": self._parse_countries_list(
                table_info.get("countries_available_for_sales", "")
            )
            or None,
            "domicile": table_info.get("domicile"),
            "management_company": table_info.get("management_company"),
            "legal_structure": table_info.get("legal_structure"),
            "isin_code": table_info.get("isin_code"),
            "stock_exchanges": table_info.get("stock_exchanges"),
            "tickers": table_info.get("tickers"),
            "dividend_policy": table_info.get("dividend_policy"),
            "minimum_initial_investment_amount": table_info.get(
                "minimum_initial_investment_amount"
            ),
            "creation_unit": table_info.get("creation_unit"),
            "initial_offering_price": table_info.get("initial_offering_price"),
            "dealing_day": table_info.get("dealing_day"),
            "dealing_deadline": table_info.get("dealing_deadline"),
            "settlement_date_subscriptions": table_info.get(
                "settlement_date_subscriptions"
            ),
            "extraction_metadata": {
                "extraction_date": datetime.now().isoformat(),
                "source_file": self.filename,
                "total_slides": len(self.prs.slides),
                "extraction_method": "improved_multi_pattern+table",
            },
        }

        extracted["performance_flags"] = self._detect_performance_content()

        self._log("\n" + "─" * 70)
        return {
            "fund_info_complete": extracted,
            "fund_info_detailed": {
                k: {"value": v, "source": "pptx"} for k, v in extracted.items()
            },
            "extraction_sources": {"pptx": extracted["extraction_metadata"]},
            "performance_flags": extracted["performance_flags"],
        }

    # ═══════════════════════════════════════════════════════════════
    # EXTRACTION ROBUSTE : INCEPTION DATE
    # ═══════════════════════════════════════════════════════════════
    def _extract_inception_date_robust(
        self, inception_from_table: Optional[str] = None
    ) -> Optional[str]:
        """
        Extraction robuste de la date de création.

        Priorité :
          1) Valeur issue de la table "General characteristics"
             (ex: '2025' dans ton PPTX)
          2) Recherche regex dans tout le document
        """
        self._log("🔎 Extraction inception_date...")

        # 1) Table "General characteristics"
        if inception_from_table:
            txt = inception_from_table.strip()
            # Cas simple : juste une année '2025'
            if re.fullmatch(r"\d{4}", txt):
                self._log(f"   ✅ Trouvé dans la table (année): {txt}")
                return txt
            # Sinon on essaie une normalisation plus fine
            normalized = self._normalize_date(txt)
            if normalized:
                self._log(f"   ✅ Trouvé dans la table (normalisé): {normalized}")
                return normalized

        # 2) Fallback : recherche dans tout le texte
        patterns = [
            r"inception\s+date[:\s]+([A-Z]{2}[/\-]\d{2}[/\-]\d{4})",
            r"inception\s+date[:\s]+(\d{2}[/\-]\d{2}[/\-]\d{4})",
            r"inception\s+date[:\s]+([A-Z][a-z]+\s+\d{4})",
            r"inception\s+date[:\s]+(\d{4})",
            r"date\s+de\s+cr[ée]ation[:\s]+(\d{2}[/\-]\d{2}[/\-]\d{4})",
        ]

        for slide_idx, slide in enumerate(self.prs.slides, 1):
            text = self._get_slide_text(slide)
            for pattern in patterns:
                matches = re.findall(pattern, text, re.IGNORECASE)
                if matches:
                    date_str = matches[0]
                    # Si c'est juste une année on la renvoie telle quelle
                    if re.fullmatch(r"\d{4}", date_str.strip()):
                        self._log(
                            f"   ✅ Trouvé dans le texte (année): {date_str.strip()}"
                        )
                        return date_str.strip()
                    normalized = self._normalize_date(date_str)
                    if normalized:
                        self._log(
                            f"   ✅ Trouvé dans le texte (normalisé): {normalized}"
                        )
                        return normalized

        self._log("   ❌ Non trouvé")
        return None

    def _normalize_date(self, date_str: str) -> Optional[str]:
        """Normalise une date en DD/MM/YYYY quand c'est possible."""
        try:
            if date_str.upper().startswith("XX"):
                date_str = "01" + date_str[2:]

            for fmt in ["%d/%m/%Y", "%d-%m-%Y", "%B %Y", "%b %Y", "%m/%Y"]:
                try:
                    dt = datetime.strptime(date_str, fmt)
                    return dt.strftime("%d/%m/%Y")
                except ValueError:
                    continue
            return None
        except Exception:
            return None

    # ═══════════════════════════════════════════════════════════════
    # EXTRACTION ROBUSTE : SRI RISK LEVEL
    # ═══════════════════════════════════════════════════════════════
    def _extract_sri_risk_level_robust(self) -> Optional[int]:
        """Extraction robuste du SRI."""
        self._log("🔎 Extraction sri_risk_level...")

        patterns = [
            r"(\d)\s*[/]\s*7",
            r"(\d)\s+(?:out\s+of|sur|over)\s+7",
            r"risk\s+(?:level|scale)[:\s]+(\d)",
        ]

        for slide in self.prs.slides:
            text = self._get_slide_text(slide)
            for pattern in patterns:
                matches = re.findall(pattern, text, re.IGNORECASE)
                if matches:
                    try:
                        level = int(matches[0])
                        if 1 <= level <= 7:
                            self._log(f"   ✅ Trouvé: niveau {level}")
                            return level
                    except ValueError:
                        continue

        self._log("   ❌ Non trouvé")
        return None

    # ═══════════════════════════════════════════════════════════════
    # EXTRACTION ROBUSTE : SFDR ARTICLE
    # ═══════════════════════════════════════════════════════════════
    def _extract_esg_sfdr_article_robust(self) -> Optional[str]:
        """Extraction robuste de l'article SFDR."""
        self._log("🔎 Extraction esg_sfdr_article...")

        patterns = [
            r"(?:sfdr\s+classification|classification\s+sfdr)[:\s]+(?:article\s+)?([689])",
            r"article\s+([689])",
        ]

        article_counts = {6: 0, 8: 0, 9: 0}

        for slide in self.prs.slides:
            text = self._get_slide_text(slide)
            for pattern in patterns:
                matches = re.findall(pattern, text, re.IGNORECASE)
                for match in matches:
                    article_num = int(match)
                    article_counts[article_num] += 1

        if sum(article_counts.values()) > 0:
            most_common = max(article_counts, key=article_counts.get)
            if article_counts[most_common] > 0:
                result = f"Article {most_common}"
                self._log(f"   ✅ Trouvé: {result}")
                return result

        self._log("   ❌ Non trouvé")
        return None

    # ═══════════════════════════════════════════════════════════════
    # EXTRACTION ROBUSTE : BENCHMARK
    # ═══════════════════════════════════════════════════════════════
    def _extract_benchmark_robust(self) -> Optional[str]:
        """Extraction robuste du benchmark."""
        self._log("🔎 Extraction benchmark...")

        known_benchmarks = ["S&P 500", "MSCI World", "MSCI Europe", "STOXX"]

        for benchmark in known_benchmarks:
            pattern = rf"{benchmark}[^,\n]{{0,50}}(?:Net|Total|Return|Index|NR|TR)"
            matches = re.findall(pattern, self.full_text, re.IGNORECASE)
            if matches:
                result = matches[0].strip()
                result = re.sub(r"\s+", " ", result)
                self._log(f"   ✅ Trouvé: {result}")
                return result

        self._log("   ❌ Non trouvé")
        return None

    def _extract_ter_robust(self) -> Optional[str]:
        """Extraction robuste du TER."""
        self._log("🔎 Extraction ter...")

        patterns = [
            r"(?:ter|total\s+expense\s+ratio)[:\s]+([0-9]+[.,]?[0-9]*\s*%)",
        ]

        for slide in self.prs.slides:
            text = self._get_slide_text(slide)
            for pattern in patterns:
                matches = re.findall(pattern, text, re.IGNORECASE)
                if matches:
                    ter = matches[0].strip().replace(",", ".")
                    self._log(f"   ✅ Trouvé: {ter}")
                    return ter

        self._log("   ❌ Non trouvé")
        return None

    def _extract_distribution_policy_robust(self) -> Optional[str]:
        """Extraction robuste de la politique de distribution."""
        self._log("🔎 Extraction distribution_policy...")

        patterns = [
            (r"\b(?:accumulation|accumulating)\b", "Accumulation"),
            (r"\b(?:distribution|distributing)\b", "Distribution"),
        ]

        for slide in self.prs.slides:
            text = self._get_slide_text(slide).lower()
            for pattern, policy in patterns:
                if re.search(pattern, text, re.IGNORECASE):
                    self._log(f"   ✅ Trouvé: {policy}")
                    return policy

        self._log("   ❌ Non trouvé")
        return None

    def _extract_fund_type_robust(self) -> Optional[str]:
        """Extraction robuste du type de fonds."""
        self._log("🔎 Extraction fund_type...")

        fund_types = [
            ("UCITS ETF", r"UCITS\s+ETF"),
            ("UCITS", r"UCITS(?!\s+ETF)"),
            ("ETF", r"\bETF\b"),
        ]

        for type_name, pattern in fund_types:
            if re.search(pattern, self.full_text, re.IGNORECASE):
                self._log(f"   ✅ Trouvé: {type_name}")
                return type_name

        self._log("   ❌ Non trouvé")
        return None

    def _extract_currency_robust(self) -> Optional[str]:
        """Extraction robuste de la devise."""
        self._log("🔎 Extraction currency...")

        patterns = [r"currency[:\s]+(USD|EUR|GBP|CHF)", r"\b(USD|EUR|GBP|CHF)\b"]

        for slide in self.prs.slides:
            text = self._get_slide_text(slide)
            for pattern in patterns:
                matches = re.findall(pattern, text, re.IGNORECASE)
                if matches:
                    currency = matches[0].upper()
                    self._log(f"   ✅ Trouvé: {currency}")
                    return currency

        self._log("   ❌ Non trouvé")
        return None

    def _extract_minimum_investment_robust(self) -> Optional[str]:
        """Extraction robuste du minimum d'investissement."""
        self._log("🔎 Extraction minimum_investment...")

        patterns = [
            r"minimum\s+(?:investment|subscription)[:\s]+([^,\n\.]{5,50})",
            r"no\s+minimum",
            r"none",
        ]

        for slide in self.prs.slides:
            text = self._get_slide_text(slide)
            for pattern in patterns:
                matches = re.findall(pattern, text, re.IGNORECASE)
                if matches:
                    result = matches[0].strip() if isinstance(matches[0], str) else ""
                    if "no" in result.lower() or "none" in result.lower():
                        result = "None"
                    self._log(f"   ✅ Trouvé: {result}")
                    return result

        self._log("   ❌ Non trouvé")
        return None

    def _extract_fund_name(self) -> Optional[str]:
        """Extrait le nom du fonds depuis la première slide."""
        if len(self.prs.slides) > 0:
            slide_1_text = self._get_slide_text(self.prs.slides[0])
            lines = [l.strip() for l in slide_1_text.split("\n") if l.strip()]
            if lines:
                return lines[0]
        return None

    def _detect_performance_content(self) -> Dict[str, Any]:
        """Détecte la présence de contenu de performance dans les slides."""
        has_perf = False
        slides_with_perf = []

        perf_keywords = ["performance", "return", "rendement", "ytd"]

        for idx, slide in enumerate(self.prs.slides, 1):
            text = self._get_slide_text(slide).lower()
            if any(kw in text for kw in perf_keywords):
                has_perf = True
                slides_with_perf.append(idx)

        return {
            "has_performance_content": has_perf,
            "slides_with_performance": slides_with_perf,
        }

    # ═══════════════════════════════════════════════════════════════
    # UTILITAIRES GÉNÉRAUX
    # ═══════════════════════════════════════════════════════════════
    def _extract_all_text(self) -> str:
        """Extrait tout le texte du PPTX (toutes slides concaténées)."""
        all_text = []
        for slide in self.prs.slides:
            all_text.append(self._get_slide_text(slide))
        return "\n".join(all_text)

    def _extract_slides_with_metadata(self) -> List[Dict[str, Any]]:
        """Extrait chaque slide avec métadonnées minimales."""
        slides_data = []
        for idx, slide in enumerate(self.prs.slides, 1):
            slides_data.append(
                {
                    "slide_number": idx,
                    "text": self._get_slide_text(slide),
                    "has_table": any(
                        s.shape_type == MSO_SHAPE_TYPE.TABLE for s in slide.shapes
                    ),
                }
            )
        return slides_data

    def _get_slide_text(self, slide) -> str:
        """Extrait le texte d'une slide (y compris contenu de tables)."""
        try:
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                text_parts.append(cell.text)
            return " ".join(text_parts)
        except Exception:
            return ""

    def _log(self, message: str):
        """Log conditionnel (activé si verbose=True)."""
        if self.verbose:
            print(message)


# ═══════════════════════════════════════════════════════════════════
# FONCTION PRINCIPALE AMÉLIORÉE
# ═══════════════════════════════════════════════════════════════════
def run_improved_extraction(
    pptx_path: str,
    prospectus_extracted: Dict[str, Any],
    output_path: str = "enriched_with_pptx_improved.json",
    llm_api_key: str = None,
    llm_base_url: str = "https://tokenfactory.esprit.tn/api",
    llm_model: str = "hosted_vllm/Llama-3.1-70B-Instruct",
) -> Dict[str, Any]:
    """
    Exécute l'extraction améliorée avec LLM.

    - Si llm_api_key est fourni → patterns + table + LLM
    - Sinon → patterns + table uniquement

    Les infos prospectus_extracted servent de fallback :
      si un champ est manquant côté PPTX, on complète avec la valeur prospectus.
    """

    print("\n" + "=" * 70)
    print("🔍 EXTRACTION PPTX AMÉLIORÉE AVEC LLM")
    print("=" * 70 + "\n")

    extractor = ImprovedPPTXExtractor(pptx_path, verbose=True)

    if llm_api_key:
        result = extractor.extract_complete_metadata_with_llm(
            api_key=llm_api_key, base_url=llm_base_url, model=llm_model
        )
    else:
        result = extractor.extract_all()

    # FUSION AVEC LES DONNÉES PROSPECTUS (FALLBACK)
    fund_info = result.get("fund_info_complete", {}) or {}
    for key, value in prospectus_extracted.items():
        if fund_info.get(key) in [None, "", "None", [], "null"]:
            fund_info[key] = value
    result["fund_info_complete"] = fund_info

    # Sauvegarde
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(result, f, indent=2, ensure_ascii=False)

    print(f"\n✅ Résultat sauvegardé: {output_path}")
    print("=" * 70 + "\n")

    return result


# ═══════════════════════════════════════════════════════════════════
# EXEMPLE D'UTILISATION (à adapter dans ton notebook Colab)
# ═══════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    PPTX_PATH = (
        "example_2/expl2.pptx"
    )
    LLM_API_KEY = "sk-7c0b80cf494746f580cc5ba555d739b2"  # ⚠️ à remplacer

    prospectus_data = {
        "fund_name": "ODDO BHF US Equity Active UCITS ETF",
        "currency": "USD",
        "ter": "0.35%",
        "distribution_policy": "Accumulation",
    }

    result = run_improved_extraction(
        pptx_path=PPTX_PATH,
        prospectus_extracted=prospectus_data,
        output_path="example_2/outputs/enriched_improved_complete_test.json",
        llm_api_key=LLM_API_KEY,
    )

    print("\n📊 MÉTADONNÉES EXTRAITES :")
    print("-" * 70)
    for field, value in result["fund_info_complete"].items():
        symbol = "✅" if value is not None else "❌"
        print(f"{symbol} {field}: {value}")