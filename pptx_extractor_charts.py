import os
from pptx import Presentation
from collections import defaultdict
from datetime import datetime, timedelta

EMU_TO_CM = 360000  # 1 cm = 360,000 EMUs

def classify_position(x, y, w, h, slide_width=33.87, slide_height=19.05):
    """
    Divide slide into 9 regions: top-left, top-center, top-right,
    middle-left, middle-center, middle-right, bottom-left, bottom-center, bottom-right.
    Middle vertical section gets bigger space.
    """
    center_x = x + w / 2
    center_y = y + h / 2

    # Vertical thresholds
    top_thresh = slide_height * 0.18
    bottom_thresh = slide_height * 0.82

    if center_y < top_thresh:
        v_region = "top"
    elif center_y > bottom_thresh:
        v_region = "bottom"
    else:
        v_region = "middle"

    # Horizontal thresholds
    left_thresh = slide_width * 0.33
    right_thresh = slide_width * 0.66

    if center_x < left_thresh:
        h_region = "left"
    elif center_x > right_thresh:
        h_region = "right"
    else:
        h_region = "center"

    semantic = f"{v_region}-{h_region}"
    return semantic


def describe_table(shape):
    """Return an LLM-friendly natural-language description of a table."""
    table = shape.table
    rows = len(table.rows)
    cols = len(table.columns)

    data = []
    for r in range(rows):
        row_data = []
        for c in range(cols):
            row_data.append(table.cell(r, c).text.strip())
        data.append(row_data)

    description = [
        f"[TABLE] {rows} rows × {cols} columns.",
        "Data:"
    ]

    for r, row in enumerate(data):
        description.append(f" Row {r+1}: " + " | ".join(row))

    return "\n".join(description)


def extract_chart_data(shape):
    """
    Extract chart metadata - optimized for speed and readability.
    Shows chart type and intelligently handles different data scenarios.
    """
    chart = shape.chart
    summary = ["[CHART] Detected chart."]

    # Chart Type
    try:
        chart_type_name = str(chart.chart_type).split('.')[-1].replace('_', ' ').title()
        summary.append(f"Type: {chart_type_name}")
    except:
        pass

    # Title
    try:
        if chart.has_title and chart.chart_title.has_text_frame:
            summary.append(f"Title: {chart.chart_title.text_frame.text}")
    except:
        pass

    # Try to get series info
    try:
        for series in chart.series:
            series_name = series.name if series.name else "Unnamed Series"
            summary.append(f"Series: {series_name}")
            
            # Get categories
            try:
                categories = list(chart.plots[0].categories)
                num_cats = len(categories)
                
                if num_cats == 0:
                    continue
                
                # Check if categories are Excel date serials (5-digit numbers 40000-50000)
                first_cat_str = str(categories[0])
                is_date_series = first_cat_str.isdigit() and 40000 <= int(first_cat_str) <= 50000
                
                # Decision: Skip detailed data for time-series with many points
                if is_date_series and num_cats > 50:
                    # Convert first and last to dates for context
                    try:
                        start_date = datetime(1899, 12, 30) + timedelta(days=int(categories[0]))
                        end_date = datetime(1899, 12, 30) + timedelta(days=int(categories[-1]))
                        summary.append(f" Time series: {start_date.strftime('%Y-%m-%d')} to {end_date.strftime('%Y-%m-%d')} ({num_cats} data points)")
                    except:
                        summary.append(f" Time series with {num_cats} data points")
                elif num_cats <= 15:
                    # Show all categories for small datasets
                    cat_list = [str(c) for c in categories]
                    summary.append(f" Categories: {', '.join(cat_list)}")
                else:
                    # Show abbreviated list for medium datasets
                    cat_list = [str(c) for c in categories[:5]]
                    summary.append(f" Categories ({num_cats} total): {', '.join(cat_list)} ... (and {num_cats - 5} more)")
                    
            except:
                pass
    except:
        pass

    return "\n".join(summary)


def extract_text_from_shape(shape, slide_width, slide_height):
    """Recursively extract all text, tables, charts, and grouped content."""
    results = []

    # ---- TABLE EXTRACTION (UPGRADED) ----
    if shape.shape_type == 19:  # TABLE
        x = shape.left / EMU_TO_CM
        y = shape.top / EMU_TO_CM
        w = shape.width / EMU_TO_CM
        h = shape.height / EMU_TO_CM
        semantic = classify_position(x, y, w, h, slide_width, slide_height)

        table_description = describe_table(shape)
        results.append((y, x, semantic, table_description))
        return results

    # ---- CHART EXTRACTION (NEW) ----
    if shape.has_chart:
        x = shape.left / EMU_TO_CM
        y = shape.top / EMU_TO_CM
        w = shape.width / EMU_TO_CM
        h = shape.height / EMU_TO_CM
        semantic = classify_position(x, y, w, h, slide_width, slide_height)

        chart_summary = extract_chart_data(shape)
        results.append((y, x, semantic, chart_summary))
        return results

    # ---- GROUP SHAPES (EXPANDED TO SUPPORT CHARTS & TABLES INSIDE GROUPS) ----
    if shape.shape_type == 6 and hasattr(shape, "shapes"):
        for sh in shape.shapes:
            results.extend(extract_text_from_shape(sh, slide_width, slide_height))
        return results

    # ---- TEXT EXTRACTION (UNCHANGED) ----
    if hasattr(shape, "text") and shape.text.strip():
        x = shape.left / EMU_TO_CM
        y = shape.top / EMU_TO_CM
        w = shape.width / EMU_TO_CM
        h = shape.height / EMU_TO_CM
        semantic = classify_position(x, y, w, h, slide_width, slide_height)
        results.append((y, x, semantic, shape.text.strip()))

    return results


def extract_pptx_to_txt(pptx_path, output_txt_path):
    prs = Presentation(pptx_path)
    slide_width = prs.slide_width / EMU_TO_CM
    slide_height = prs.slide_height / EMU_TO_CM

    lines = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        lines.append(f"=== Slide {slide_idx} ===\n")

        all_items = []
        for shape in slide.shapes:
            all_items.extend(extract_text_from_shape(shape, slide_width, slide_height))

        all_items_sorted = sorted(all_items, key=lambda t: (t[0], t[1]))

        semantic_dict = defaultdict(list)
        for y, x, semantic, content in all_items_sorted:
            semantic_dict[semantic].append((y, x, content))

        for semantic, items in semantic_dict.items():
            lines.append(f"[{semantic}]:\n")
            for _, _, content in sorted(items, key=lambda t: (t[0], t[1])):
                lines.append(f"{content}\n")
        lines.append("\n")

    with open(output_txt_path, "w", encoding="utf-8") as f:
        f.writelines(lines)

    print(f"✅ Extracted text saved to {output_txt_path}")


# === RUN - PROCESSES ALL PPTX FILES ===
if __name__ == "__main__":
    folder = "example_3"
    # Filter out temporary PowerPoint files (starting with ~$)
    pptx_files = [f for f in os.listdir(folder) if f.endswith(".pptx") and not f.startswith("~$")]
    
    if not pptx_files:
        raise FileNotFoundError("No .pptx file found in the folder.")
    
    print(f"\n🔍 Found {len(pptx_files)} PPTX file(s) in '{folder}'")
    print("="*60)
    
    # Process EVERY PPTX FILE
    for idx, pptx_file in enumerate(pptx_files, 1):
        pptx_path = os.path.join(folder, pptx_file)
        
        # Create unique output filename for each pptx
        base_name = os.path.splitext(pptx_file)[0]
        output_txt = os.path.join(folder, f"slides_extracted_{base_name}.txt")
        
        print(f"\n📄 [{idx}/{len(pptx_files)}] Processing: {pptx_file}")
        extract_pptx_to_txt(pptx_path, output_txt)
    
    print("\n" + "="*60)
    print(f"✅ Successfully processed {len(pptx_files)} PPTX file(s)!")
    print(f"📁 Output files saved in: {folder}/")