import os
from pptx import Presentation
from collections import defaultdict

EMU_TO_CM = 360000  # 1 cm = 360,000 EMUs

def classify_position(x, y, w, h, slide_width=33.87, slide_height=19.05):
    """
    Divide slide into 9 regions: top-left, top-center, top-right,
    middle-left, middle-center, middle-right, bottom-left, bottom-center, bottom-right.
    Middle vertical section gets bigger space.
    """
    center_x = x + w / 2
    center_y = y + h / 2

    # Vertical thresholds
    top_thresh = slide_height * 0.18       # top 25%
    bottom_thresh = slide_height * 0.82    # bottom 25%, middle 50%

    if center_y < top_thresh:
        v_region = "top"
    elif center_y > bottom_thresh:
        v_region = "bottom"
    else:
        v_region = "middle"

    # Horizontal thresholds
    left_thresh = slide_width * 0.33
    right_thresh = slide_width * 0.66

    if center_x < left_thresh:
        h_region = "left"
    elif center_x > right_thresh:
        h_region = "right"
    else:
        h_region = "center"

    semantic = f"{v_region}-{h_region}"
    return semantic

def extract_text_from_shape(shape, slide_width, slide_height):
    """Recursively extract all text from a shape (including groups and tables)."""
    texts = []

    # Table
    if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
        x = shape.left / EMU_TO_CM
        y = shape.top / EMU_TO_CM
        w = shape.width / EMU_TO_CM
        h = shape.height / EMU_TO_CM
        semantic = classify_position(x, y, w, h, slide_width, slide_height)

        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    texts.append((y, x, semantic, cell.text.strip()))

    # Group
    elif shape.shape_type == 6 and hasattr(shape, "shapes"):  # GROUP
        for sh in shape.shapes:
            texts.extend(extract_text_from_shape(sh, slide_width, slide_height))

    # Regular text box / placeholder
    elif hasattr(shape, "text") and shape.text.strip():
        x = shape.left / EMU_TO_CM
        y = shape.top / EMU_TO_CM
        w = shape.width / EMU_TO_CM
        h = shape.height / EMU_TO_CM
        semantic = classify_position(x, y, w, h, slide_width, slide_height)
        texts.append((y, x, semantic, shape.text.strip()))

    return texts

def extract_pptx_to_txt(pptx_path, output_txt_path):
    prs = Presentation(pptx_path)
    slide_width = prs.slide_width / EMU_TO_CM
    slide_height = prs.slide_height / EMU_TO_CM

    lines = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        lines.append(f"=== Slide {slide_idx} ===\n")

        # Extract all text shapes
        all_texts = []
        for shape in slide.shapes:
            all_texts.extend(extract_text_from_shape(shape, slide_width, slide_height))

        # Sort: top -> bottom, left -> right
        all_texts_sorted = sorted(all_texts, key=lambda t: (t[0], t[1]))

        # Group texts by semantic
        semantic_dict = defaultdict(list)
        for y, x, semantic, text in all_texts_sorted:
            semantic_dict[semantic].append((y, x, text))

        # Print each semantic once with all texts in y->x order
        for semantic, items in semantic_dict.items():
            lines.append(f"[{semantic}]:\n")
            for _, _, text in sorted(items, key=lambda t: (t[0], t[1])):
                lines.append(f"{text}\n")

        lines.append("\n")  # blank line between slides

    # Save
    with open(output_txt_path, "w", encoding="utf-8") as f:
        f.writelines(lines)

    print(f"✅ Extracted text saved to {output_txt_path}")


# === Run Example ===
if __name__ == "__main__":
    folder = "example_2"
    pptx_files = [f for f in os.listdir(folder) if f.endswith(".pptx")]
    if not pptx_files:
        raise FileNotFoundError("No .pptx file found in the folder.")

    pptx_path = os.path.join(folder, pptx_files[0])
    output_txt = os.path.join(folder, "slides_extracted_text.txt")

    extract_pptx_to_txt(pptx_path, output_txt)
